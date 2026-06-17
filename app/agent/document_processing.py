import io
import uuid
from pathlib import Path
from typing import List, Dict, Any, Tuple

from PIL import Image
from loguru import logger

from app.utils.ocr_utils import extract_text_from_image, ocr_is_available
from app.utils.vision_utils import analyze_image


def _process_image(image: Image.Image, source: str, page: int) -> List[Dict[str, Any]]:
    """Run both Vision API and OCR on an image; return document dicts."""
    docs = []

    # GPT-4o Vision analysis
    vision_text = analyze_image(image)
    if vision_text.strip():
        docs.append({
            "text": vision_text,
            "source": source,
            "page": page,
            "chunk_type": "vision",
        })

    # Tesseract OCR as fallback / supplement
    if ocr_is_available():
        ocr_text = extract_text_from_image(image)
        if ocr_text.strip() and ocr_text != vision_text:
            docs.append({
                "text": ocr_text,
                "source": source,
                "page": page,
                "chunk_type": "ocr",
            })

    return docs


def load_pdf(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Extract text and images from a PDF.
    Returns list of {text, source, page, chunk_type} dicts.
    """
    import fitz  # PyMuPDF

    documents = []
    pdf = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num, page in enumerate(pdf):
        # Extract text
        text = page.get_text("text").strip()
        if text:
            documents.append({
                "text": text,
                "source": filename,
                "page": page_num + 1,
                "chunk_type": "text",
            })

        # Extract images
        for img_index, img_ref in enumerate(page.get_images(full=True)):
            xref = img_ref[0]
            base_image = pdf.extract_image(xref)
            image_bytes = base_image["image"]
            try:
                pil_image = Image.open(io.BytesIO(image_bytes))
                image_docs = _process_image(pil_image, filename, page_num + 1)
                documents.extend(image_docs)
            except Exception as e:
                logger.warning(f"Failed to process image on page {page_num + 1}: {e}")

    pdf.close()
    logger.info(f"PDF '{filename}': extracted {len(documents)} document sections.")
    return documents


def load_pptx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Extract text and images from a PPTX presentation.
    Returns list of {text, source, page, chunk_type} dicts.
    """
    from pptx import Presentation
    from pptx.util import Inches

    documents = []
    prs = Presentation(io.BytesIO(file_bytes))

    for slide_num, slide in enumerate(prs.slides):
        slide_text_parts = []

        for shape in slide.shapes:
            # Text extraction
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in paragraph.runs).strip()
                    if line:
                        slide_text_parts.append(line)

            # Image extraction
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_bytes = shape.image.blob
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    image_docs = _process_image(pil_image, filename, slide_num + 1)
                    documents.extend(image_docs)
                except Exception as e:
                    logger.warning(f"Failed to process image on slide {slide_num + 1}: {e}")

        if slide_text_parts:
            documents.append({
                "text": "\n".join(slide_text_parts),
                "source": filename,
                "page": slide_num + 1,
                "chunk_type": "text",
            })

    logger.info(f"PPTX '{filename}': extracted {len(documents)} document sections.")
    return documents


def load_document(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Route to appropriate loader based on file extension."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return load_pdf(file_bytes, filename)
    elif ext in (".pptx", ".ppt"):
        return load_pptx(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .pptx")
